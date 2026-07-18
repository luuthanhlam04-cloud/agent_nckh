"""
parser.py - Tầng Bóc tách Tài liệu Học thuật (Document Parsing Layer)
======================================================================
Hỗ trợ:
  - PDFParser  : Sử dụng PyMuPDF (fitz) bóc văn bản -> Markdown theo cấp đoạn văn.
  - PPTXParser : Trích xuất text từ shapes, sau đó gửi lên Gemini Flash để làm sạch
                 và cấu trúc hóa nội dung học thuật. Gài time.sleep(4) giữa mỗi slide
                 để chống Rate Limit Google (15 req/phút).
  - chunk_by_paragraph() : Chiến lược cắt chunk theo ranh giới ngữ nghĩa của đoạn văn.

Thiết kế:
  - Mỗi parser trả về list[dict] chuẩn, sẵn sàng đẩy vào QdrantManager.upsert_chunks().
  - Chunk size và overlap có thể điều chỉnh qua hằng số cấu hình.
  - Xử lý lỗi chi tiết: ghi log cảnh báo và bỏ qua slide lỗi thay vì crash.

Ghi chú PPTXParser:
  Hiện tại sử dụng chế độ Text+Gemini (trích xuất text từ shapes -> đẩy lên Gemini
  để làm sạch và tóm tắt học thuật). Chế độ Vision Image thật sự (render slide
  thành PNG bằng win32com) là nâng cấp tương lai — cần Microsoft Office cài sẵn.
"""

import os
import time
import logging
import re
import uuid
from pathlib import Path
from typing import List, Dict, Any, Optional

import fitz  # PyMuPDF
import google.genai as genai
from google.genai import types as genai_types

# python-pptx chỉ dùng để đọc slide text và (tương lai) xuất slide thành ảnh
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    Presentation = None  # type: ignore

# ─── Logging ──────────────────────────────────────────────────────────────────
# [S4-FIX] Không gọi basicConfig ở đây — main.py đã cấu hình toàn cục với
# FileHandler + StreamHandler. Gọi lại chỉ tạo duplicate handler.
logger = logging.getLogger("Parser")

# ─── Config (đọc từ env đã được load_dotenv() trong main.py) ──────────────────
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY", "")

# ─── Hằng số cấu hình Chunking ────────────────────────────────────────────────
# [Sprint A] Giảm chunk size để vector đặc trưng hơn, tránh Semantic Dilution.
# MAX 600 ≈ 150 từ = 1-2 đoạn văn → mỗi vector đại diện cho 1 ý cụ thể.
MIN_CHUNK_CHARS = 80        # Không bỏ qua bullet point, định nghĩa ngắn quan trọng
MAX_CHUNK_CHARS = 600       # Giới hạn ~150 từ để vector đặc trưng (giảm từ 2000)
CHUNK_OVERLAP_CHARS = 150   # ~1 câu cuối để giữ ngữ cảnh ranh giới (giảm từ 400)
PPTX_SLEEP_SECONDS = 4      # Ngủ giữa mỗi slide để tránh Rate Limit Google (15 req/phút)


# ══════════════════════════════════════════════════════════════════════════════
#  HELPER: Chunking theo ranh giới đoạn văn (Recursive + Heading-aware)
# ══════════════════════════════════════════════════════════════════════════════
def chunk_by_paragraph(
    text: str,
    source: str,
    page: int = 0,
    metadata: Optional[Dict] = None,
    section_title: str = "",
) -> List[Dict[str, Any]]:
    """
    [Sprint A] Cắt văn bản thành chunk nhỏ theo ranh giới ngữ nghĩa.
    - Ưu tiên tách tại: heading markdown (# ## ###) → dòng trống → câu.
    - Mỗi chunk mang metadata section_title để AI biết ngữ cảnh Section.

    Args:
        text         : Văn bản đầu vào.
        source       : Tên file nguồn (lưu vào Qdrant payload).
        page         : Số trang PDF.
        metadata     : Thông tin bổ sung.
        section_title: Tên Section hiện tại (Abstract, Method, Result...)

    Returns:
        Danh sách dict chunk chuẩn để đưa vào QdrantManager.
    """
    # Tách theo dòng trống ≥ 2 dòng, hoặc ngay trước heading Markdown (# ## ###)
    # [Sprint A] Heading tạo ranh giới cứng → mỗi Section thành cụm chunk riêng
    paragraphs = re.split(r"\n{2,}|(?=^#{1,3}\s)", text, flags=re.MULTILINE)
    chunks = []
    buffer = ""
    current_section = section_title  # Theo dõi Section hiện tại

    for para in paragraphs:
        para = para.strip()
        if not para:
            continue

        # Nhận diện heading Markdown → flush buffer và cập nhật section_title
        # [Sprint A] Heading = ranh giới cứng, không ghép vào chunk đang dở
        heading_match = re.match(r'^(#{1,3})\s+(.+?)(?:\n|$)', para)
        if heading_match:
            # Flush buffer hiện tại trước khi bắt đầu Section mới
            if buffer and len(buffer) >= MIN_CHUNK_CHARS:
                chunks.append({
                    "text": buffer,
                    "source": source,
                    "page": page,
                    "metadata": {**(metadata or {}), "section_title": current_section},
                })
            buffer = ""
            current_section = heading_match.group(2).strip()
            # Bóc phần text còn lại (nếu có) nằm ngay dưới heading mà bị dính liền
            para = para[heading_match.end():].strip()
            if not para:
                continue

        # Nếu buffer + đoạn hiện tại nằm trong giới hạn → ghép tiếp
        if len(buffer) + len(para) + 2 < MAX_CHUNK_CHARS:
            buffer = (buffer + "\n\n" + para).strip() if buffer else para
        else:
            # Buffer đầy → flush ra chunk
            if len(buffer) >= MIN_CHUNK_CHARS:
                chunks.append({
                    "text": buffer,
                    "source": source,
                    "page": page,
                    "metadata": {**(metadata or {}), "section_title": current_section},
                })
                # Overlap: giữ lại ~1 câu cuối để không mất ngữ cảnh ranh giới
                if len(buffer) > CHUNK_OVERLAP_CHARS:
                    overlap_text = buffer[-CHUNK_OVERLAP_CHARS:]
                    # Cắt tại khoảng trắng gần nhất để không đứt ngang từ
                    space_idx = overlap_text.find(' ')
                    if space_idx != -1 and space_idx < len(overlap_text) - 1:
                        overlap_text = overlap_text[space_idx + 1:]
                    buffer = (overlap_text + "\n\n" + para).strip()
                else:
                    buffer = para
            else:
                # Buffer quá ngắn → ghép tiếp thay vì flush rác
                buffer = (buffer + "\n\n" + para).strip() if buffer else para

    # Flush phần còn lại
    if buffer and len(buffer) >= MIN_CHUNK_CHARS:
        chunks.append({
            "text": buffer,
            "source": source,
            "page": page,
            "metadata": {**(metadata or {}), "section_title": current_section},
        })

    logger.info(
        "[Parser] Chunking '%s' trang %d: %d chunks (section='%s')",
        source, page, len(chunks), current_section or "unknown"
    )
    return chunks


def _extract_blocks_with_headings(page: "fitz.Page", file_name: str) -> List[Dict]:
    """
    [Sprint A] Đọc trang PDF theo chế độ 'dict' để lấy thêm font metadata.
    Phát hiện heading dựa vào font_size > median * 1.2 hoặc is_bold.
    Trả về danh sách blocks có đánh dấu is_heading, text, và font_size.

    Fallback an toàn: nếu get_text('dict') thất bại → dùng get_text('text') bình thường.
    """
    blocks_out = []
    try:
        raw_dict = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE)
        all_sizes = []
        # Thu thập tất cả font_size để tính median
        for block in raw_dict.get("blocks", []):
            if block.get("type") != 0:  # type=0 là text block
                continue
            for line in block.get("lines", []):
                for span in line.get("spans", []):
                    sz = span.get("size", 0)
                    if sz > 0:
                        all_sizes.append(sz)

        if not all_sizes:
            # Fallback: không có font info, dùng plain text
            return [{"text": page.get_text("text"), "is_heading": False, "font_size": 0}]

        # Median font size = ngưỡng phân biệt body text vs heading
        all_sizes_sorted = sorted(all_sizes)
        median_size = all_sizes_sorted[len(all_sizes_sorted) // 2]
        heading_threshold = median_size * 1.2  # Dòng nào to hơn 120% median → heading

        # Nhóm spans thành blocks văn bản, gắn cờ is_heading
        for block in raw_dict.get("blocks", []):
            if block.get("type") != 0:
                continue
            block_lines = []
            block_sizes = []
            for line in block.get("lines", []):
                line_text = ""
                line_size = 0
                for span in line.get("spans", []):
                    span_text = span.get("text", "").strip()
                    if span_text:
                        line_text += span_text + " "
                        line_size = max(line_size, span.get("size", 0))
                line_text = line_text.strip()
                if len(line_text) > 5:  # Bỏ qua dòng quá ngắn (số trang...)
                    block_lines.append(line_text)
                    block_sizes.append(line_size)

            if not block_lines:
                continue

            block_text = " ".join(block_lines)
            avg_size = sum(block_sizes) / len(block_sizes) if block_sizes else 0
            is_heading = avg_size >= heading_threshold and len(block_text) < 200

            blocks_out.append({
                "text": block_text,
                "is_heading": is_heading,
                "font_size": round(avg_size, 1),
            })

    except Exception as e:
        logger.warning("[PDFParser] Fallback plain text cho trang (dict mode lỗi): %s", e)
        blocks_out = [{"text": page.get_text("text"), "is_heading": False, "font_size": 0}]

    return blocks_out


# ══════════════════════════════════════════════════════════════════════════════
#  PDFParser - Bóc tách PDF bằng PyMuPDF
# ══════════════════════════════════════════════════════════════════════════════
class PDFParser:
    """
    Bóc tách tài liệu PDF thành các chunk văn bản chuẩn.
    - Sử dụng PyMuPDF (fitz) đọc từng trang.
    - Lọc các dòng rác (số trang, tiêu đề lặp, dòng quá ngắn).
    - Kết xuất Markdown đơn giản (heading, paragraph) trước khi chunk.
    """

    def parse(self, pdf_path: str) -> List[Dict[str, Any]]:
        """
        Bóc tách toàn bộ file PDF.

        Args:
            pdf_path: Đường dẫn tuyệt đối đến file PDF.

        Returns:
            Danh sách chunk chuẩn, sẵn sàng đưa vào Qdrant.
        """
        file_name = Path(pdf_path).name
        logger.info("[PDFParser] Bắt đầu bóc tách (Heading-aware): %s", file_name)

        all_chunks = []
        try:
            doc = fitz.open(pdf_path)
            current_section = ""  # [Sprint A] Theo dõi Section xuyên suốt toàn tài liệu

            for page_num, page in enumerate(doc, start=1):
                # [Sprint A] Dùng dict mode để detect heading qua font_size
                blocks = _extract_blocks_with_headings(page, file_name)

                # Gộp các block thành văn bản, đánh dấu heading bằng prefix ##
                page_lines = []
                for block in blocks:
                    text = block["text"].strip()
                    if not text or len(text) < 10:
                        continue
                    if block["is_heading"]:
                        # Prefix ## để chunk_by_paragraph nhận ra heading và flush
                        page_lines.append(f"## {text}")
                        current_section = text  # Cập nhật Section tracker
                    else:
                        page_lines.append(text)

                cleaned_text = "\n\n".join(page_lines)
                if not cleaned_text or len(cleaned_text.strip()) < 50:
                    continue

                # [Sprint B] Tạo Chunk Cha (Parent) đại diện cho toàn bộ trang/section
                parent_id = str(uuid.uuid4())
                parent_chunk = {
                    "text": cleaned_text,
                    "source": file_name,
                    "page": page_num,
                    "metadata": {
                        "file_type": "pdf",
                        "total_pages": len(doc),
                        "section_title": current_section,
                        "chunk_type": "parent",
                        "chunk_id": parent_id
                    }
                }
                all_chunks.append(parent_chunk)

                # [Sprint B] Tạo Chunk Con (Child) và liên kết bằng parent_id
                page_chunks = chunk_by_paragraph(
                    text=cleaned_text,
                    source=file_name,
                    page=page_num,
                    metadata={
                        "file_type": "pdf", 
                        "total_pages": len(doc),
                        "chunk_type": "child",
                        "parent_id": parent_id
                    },
                    section_title=current_section,
                )
                all_chunks.extend(page_chunks)

                # Cập nhật current_section từ các heading được detect trong trang này
                for block in blocks:
                    if block["is_heading"] and block["text"].strip():
                        current_section = block["text"].strip()

            doc.close()

        except Exception as e:
            logger.error("[PDFParser] Lỗi khi bóc tách '%s': %s", file_name, e)

        logger.info("[PDFParser] Hoàn thành: %s -> %d chunks.", file_name, len(all_chunks))
        return all_chunks

    def save_markdown(self, pdf_path: str, output_dir: str) -> str:
        """
        Bóc tách PDF và lưu kết quả ra file Markdown trong thư mục 02_Knowledge/.

        Args:
            pdf_path  : Đường dẫn file PDF nguồn.
            output_dir: Thư mục đích (Obsidian_Vault/02_Knowledge/).

        Returns:
            Đường dẫn file Markdown đã lưu.
        """
        file_stem = Path(pdf_path).stem
        os.makedirs(output_dir, exist_ok=True)
        output_path = os.path.join(output_dir, f"{file_stem}.md")

        chunks = self.parse(pdf_path)
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(f"# {file_stem}\n\n")
            for chunk in chunks:
                f.write(f"<!-- Page {chunk['page']} -->\n")
                f.write(chunk["text"] + "\n\n---\n\n")

        logger.info(f"[PDFParser] Đã lưu Markdown: {output_path}")
        return output_path


# ══════════════════════════════════════════════════════════════════════════════
#  PPTXParser - Bóc tách Slide PPTX bằng Gemini Vision API
# ══════════════════════════════════════════════════════════════════════════════
class PPTXParser:
    """
    Bóc tách file Slide PowerPoint theo 2 chế độ:

    Chế độ hiện tại - Text+Gemini (mặc định):
      1. Trích xuất text thô từ từng shape trong slide bằng python-pptx.
      2. Gửi text lên Gemini Flash để làm sạch cấu trúc, nhận diện thuật ngữ học thuật.
      3. Gài time.sleep(4) giữa mỗi slide để chống Rate Limit (15 req/phút Google Free).

    Chế độ tương lai - Vision Image (nâng cấp sau):
      Render slide thành PNG qua win32com (yêu cầu Microsoft Office cài sẵn).
      Gửi ảnh thật lên Gemini Vision API để đọc biểu đồ, hình minh họa.
      Sử dụng USE_VISION_MODE = True khi có Microsoft Office.

    Yêu cầu: GEMINI_API_KEY trong file .env.
    """

    # Cờ cho phép chuyển sang chế độ Vision thật sự (nâng cấp tương lai)
    USE_VISION_MODE = False

    GEMINI_PROMPT_TEXT = """Bạn là trợ lý học thuật. Dưới đây là nội dung text thô trích xuất từ một slide thuyết trình.
Nhiệm vụ: Làm sạch, cấu trúc hóa và giữ lại toàn bộ nội dung học thuật quan trọng.
Trả về văn bản thuần túy (plain text), không dùng JSON, không thêm bình luận."""

    def __init__(self):
        self._client = None
        if not GEMINI_API_KEY or "điền" in GEMINI_API_KEY.lower():
            logger.warning(
                "[PPTXParser] Thiếu GEMINI_API_KEY. Sẽ chỉ dùng text thô từ shapes."
            )
        else:
            self._client = genai.Client(api_key=GEMINI_API_KEY)

    # ── Hook cho chế độ Vision Image tương lai ─────────────────────────────────
    # Để bật chế độ Vision:
    #   1. Đặt USE_VISION_MODE = True
    #   2. Thêm method _render_slide_to_png(slide) dùng win32com.client
    #   3. Gửi kết quả PNG dạng bytes lên Gemini với mime_type="image/png"

    def parse(self, pptx_path: str) -> List[Dict[str, Any]]:
        """
        Bóc tách toàn bộ file PPTX.

        Args:
            pptx_path: Đường dẫn tuyệt đối đến file PPTX.

        Returns:
            Danh sách chunk chuẩn, sẵn sàng đưa vào Qdrant.
        """
        if not PPTX_AVAILABLE:
            logger.error("[PPTXParser] Thiếu thư viện python-pptx hoặc Pillow.")
            return []

        file_name = Path(pptx_path).name
        logger.info(f"[PPTXParser] Bắt đầu bóc tách: {file_name}")

        prs = Presentation(pptx_path)
        all_chunks = []
        total_slides = len(prs.slides)

        for slide_idx, slide in enumerate(prs.slides, start=1):
            logger.info(f"[PPTXParser] Đang xử lý slide {slide_idx}/{total_slides}...")

            # Bước 1: Luôn trích xuất text thô từ shapes (nhanh, không cần API)
            raw_texts = [
                shape.text.strip()
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            ]
            raw_slide_text = "\n".join(raw_texts)

            if not raw_slide_text.strip():
                logger.info(f"[PPTXParser] Slide {slide_idx} rỗng, bỏ qua.")
                # Vẫn ngủ để giữ kịoảng cách giữa các API call
                if self._client and slide_idx < total_slides:
                    time.sleep(PPTX_SLEEP_SECONDS)
                continue

            slide_text = raw_slide_text

            # Bước 2: Nếu có Gemini client, gửi text lên để làm sạch cấu trúc học thuật
            # (Thực sự gọi API - không bỏ qua)
            if self._client:
                try:
                    prompt = (
                        self.GEMINI_PROMPT_TEXT +
                        f"\n\nNỘI DUNG SLIDE {slide_idx}:\n{raw_slide_text}"
                    )
                    response = self._client.models.generate_content(
                        model="gemini-1.5-flash",
                        contents=prompt,
                        config=genai_types.GenerateContentConfig(
                            temperature=0.1,
                            max_output_tokens=1024,
                        ),
                    )
                    slide_text = response.text.strip()
                    logger.info(
                        f"[PPTXParser] Gemini đã làm sạch slide {slide_idx}: "
                        f"{len(raw_slide_text)} -> {len(slide_text)} ky tu"
                    )
                except Exception as e:
                    logger.warning(
                        f"[PPTXParser] Gemini lỗi slide {slide_idx}: {e}. "
                        f"Giữ lại text thô."
                    )
                    slide_text = raw_slide_text  # Fallback về text thô nếu Gemini lỗi

            # Chunk nội dung slide vừa trích xuất
            if slide_text and len(slide_text.strip()) >= 30:
                slide_chunks = chunk_by_paragraph(
                    text=slide_text,
                    source=file_name,
                    page=slide_idx,
                    metadata={"file_type": "pptx", "total_slides": total_slides},
                )
                all_chunks.extend(slide_chunks)

            # *** CHỐT CHẶN AN TOÀN: Ngủ 4 giây giữa mỗi slide ***
            # Tránh vượt quá giới hạn 15 requests/phút của Google AI Free Tier
            # [BUG-1 FIX] Dùng self._client (genai.Client) thay vì self._model (không tồn tại)
            if self._client and slide_idx < total_slides:
                logger.info(f"[PPTXParser] Đang chờ {PPTX_SLEEP_SECONDS}s (chống Rate Limit)...")
                time.sleep(PPTX_SLEEP_SECONDS)

        logger.info(f"[PPTXParser] Hoàn thành: {file_name} -> {len(all_chunks)} chunks.")
        return all_chunks

    def save_markdown(self, pptx_path: str, output_dir: str) -> str:
        """
        Bóc tách PPTX và lưu kết quả ra file Markdown trong thư mục 02_Knowledge/.
        """
        file_stem = Path(pptx_path).stem
        os.makedirs(output_dir, exist_ok=True)
        output_path = os.path.join(output_dir, f"{file_stem}.md")

        chunks = self.parse(pptx_path)
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(f"# {file_stem}\n\n")
            for chunk in chunks:
                f.write(f"<!-- Slide {chunk['page']} -->\n")
                f.write(chunk["text"] + "\n\n---\n\n")

        logger.info(f"[PPTXParser] Đã lưu Markdown: {output_path}")
        return output_path


# ─── Factory function tiện lợi ────────────────────────────────────────────────
def parse_document(file_path: str) -> List[Dict[str, Any]]:
    """
    Auto-detect loại tài liệu và gọi parser tương ứng.
    Trả về danh sách chunk chuẩn.

    Args:
        file_path: Đường dẫn đến file PDF hoặc PPTX.

    Returns:
        Danh sách chunk hoặc list rỗng nếu định dạng không hỗ trợ.
    """
    ext = Path(file_path).suffix.lower()
    if ext == ".pdf":
        return PDFParser().parse(file_path)
    elif ext in (".pptx", ".ppt"):
        return PPTXParser().parse(file_path)
    else:
        logger.warning(f"[Parser] Định dạng không hỗ trợ: '{ext}'. Bỏ qua file.")
        return []


# ─── Test nhanh khi chạy trực tiếp ───────────────────────────────────────────
if __name__ == "__main__":
    from dotenv import load_dotenv
    load_dotenv()  # Chỉ load khi chạy file độc lập để test

    import sys
    if len(sys.argv) < 2:
        print("Cách dùng: python parser.py <đường_dẫn_file.pdf|.pptx>")
        sys.exit(1)

    test_path = sys.argv[1]
    print(f"\n--- Test Parser: {test_path} ---")
    chunks = parse_document(test_path)
    print(f"\nKết quả: {len(chunks)} chunks")
    for i, c in enumerate(chunks[:3]):
        print(f"\n[Chunk {i+1}] Source: {c['source']} | Page: {c['page']}")
        print(f"  Text: {c['text'][:150]}...")
